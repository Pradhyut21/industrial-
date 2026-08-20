"""
PPTX (PowerPoint) ingestion module for the Continuity Vault.

Parses .pptx files using python-pptx, extracts slide text and speaker notes,
then flows through the existing ingest_document() RAG pipeline.

Requires: python-pptx (added to requirements.txt)
"""
from __future__ import annotations

import datetime
import io
from typing import Optional

from backend.database import get_db_connection
from backend.ingestion import ingest_document

try:
    from pptx import Presentation
    from pptx.util import Pt

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print(
        "[PptxIngestion] WARNING: python-pptx not installed. "
        "PPTX ingestion will use stub mode. Run: pip install python-pptx"
    )


def _extract_pptx_text(file_bytes: bytes) -> str:
    """
    Extracts all slide text + speaker notes from a .pptx file.
    Returns a single plain-text string suitable for indexing.
    """
    if not PPTX_AVAILABLE:
        return "[STUB] python-pptx not installed — raw PPTX text extraction unavailable."

    prs = Presentation(io.BytesIO(file_bytes))
    sections = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text.strip())

        notes_text = ""
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                notes_text = f"Speaker notes: {notes}"

        section = f"Slide {i}:\n" + "\n".join(slide_texts)
        if notes_text:
            section += f"\n{notes_text}"
        sections.append(section)

    return "\n\n".join(sections)


def ingest_pptx(
    person_id: int,
    person_name: str,
    filename: str,
    file_bytes: bytes,
    sensitivity_level: str,
    domain: str,
) -> dict:
    """
    Parses a PPTX file, extracts text, indexes into RAG pipeline,
    and creates a vault_artifact record.

    Returns a dict describing the created artifact.
    """
    raw_content = _extract_pptx_text(file_bytes)

    # Brief summary from LLM (or template fallback)
    summary = _summarise_pptx(filename, raw_content, person_name)

    # Index into RAG pipeline
    rag_doc = ingest_document(
        title=f"[PPTX] {person_name}: {filename}",
        content=raw_content + "\n\nSummary: " + summary,
        doc_type="Presentation",
        forced_author=person_name,
    )

    conn = get_db_connection()
    cursor = conn.cursor()
    now = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")

    cursor.execute(
        """
        INSERT INTO vault_artifacts
            (person_id, artifact_type, source_ref, raw_content, plain_language_summary,
             sensitivity_level, domain, ingested_at, doc_id)
        VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)
        """,
        (
            person_id,
            "pptx",
            filename,
            raw_content[:10000],  # cap storage; full text is in RAG index
            summary,
            sensitivity_level,
            domain,
            now,
            rag_doc.get("id"),
        ),
    )
    conn.commit()
    artifact_id = cursor.lastrowid
    conn.close()

    return {
        "artifact_id": artifact_id,
        "artifact_type": "pptx",
        "plain_language_summary": summary,
        "doc_id": rag_doc.get("id"),
    }


def _summarise_pptx(filename: str, text: str, author: str) -> str:
    """Summarise presentation content via Groq or template fallback."""
    import os
    from backend.llm import get_groq_response, APIConfig

    live_key = os.environ.get("GROQ_API_KEY", "") or APIConfig.key
    if live_key and text.strip():
        try:
            system = (
                "You are a technical writer summarising a presentation slide deck for a non-technical "
                "reader in an industrial organisation. In 2-4 plain-English sentences, describe: "
                "what topics this deck covers, any key action items or open issues mentioned, "
                "and why it is operationally important. Do not use jargon without explanation."
            )
            prompt = f"Presentation: {filename}\nAuthor: {author}\n\nContent:\n{text[:3000]}"
            return get_groq_response(prompt, system).strip()
        except Exception as e:
            print(f"[PptxIngestion] LLM summarise failed: {e}")

    # Template fallback
    slide_count = text.count("Slide ")
    return (
        f"Presentation '{filename}' by {author} containing approximately {slide_count} slides. "
        "Topics extracted from slide text and speaker notes. "
        "Full content has been indexed for retrieval."
    )
